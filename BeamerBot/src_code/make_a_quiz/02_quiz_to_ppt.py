"""
This script reads quiz questions and answers from an Excel file generated by make_a_quiz.py and converts them into a
PowerPoint presentation. Each slide in the presentation contains a question followed by the corresponding answer
on a separate slide.

Workflow:
1. **Load Quiz Data**: Reads the quiz questions and answers from an Excel file, including multiple-choice
   options and the correct answer for each question.
2. **PowerPoint Generation**:
   - For each question, a slide is created with the question text and answer choices.
   - A corresponding answer slide is generated with the correct answer displayed.
   - Customizable slide backgrounds and font sizes are applied to ensure readability.
3. **Saving the Output**: The generated PowerPoint presentation is saved as a `.pptx` file for review or presentation in class.

Dependencies:
- Requires the `python-pptx` library for creating and formatting the PowerPoint slides.
- The Excel file with quiz data should be formatted with columns for 'question', 'A)', 'B)', 'C)', 'D)', and 'correct_answer'.
- Ensure the necessary environment variables (e.g., `syllabus_path`) are set correctly for the input and output directories.

This script outputs a PowerPoint presentation with each question and answer in a consistent slide format, ready to be presented.
"""


# base libraries
import os
from pathlib import Path

import pandas as pd
# env setup
from dotenv import load_dotenv
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

load_dotenv()

OPENAI_KEY = os.getenv('openai_key')
OPENAI_ORG = os.getenv('openai_org')

# Path definitions
readingDir = Path(os.getenv('readingsDir'))
slideDir = Path(os.getenv('slideDir'))
syllabus_path = Path(os.getenv('syllabus_path'))

inputDir = syllabus_path.parent / "midterm"

# %%

# Load the quiz content from the Excel file
df = pd.read_excel(inputDir / "review_quiz.xlsx")


# %%
# Create a new PowerPoint presentation
prs = Presentation()

# Function to add a slide with a title and content


def add_slide(prs, title, content, answer=None, is_answer=False, bg_color=(255, 255, 255)):
    slide_layout = prs.slide_layouts[1]  # Layout with title and content
    slide = prs.slides.add_slide(slide_layout)

    # Set slide background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(bg_color[0], bg_color[1], bg_color[2])

    # Set title text and font size
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_text_frame = title_placeholder.text_frame
    title_text_frame.paragraphs[0].font.size = Pt(24)  # Set title font size

    # Set content text and font size
    text_placeholder = slide.shapes.placeholders[1]
    text_placeholder.text = content
    content_text_frame = text_placeholder.text_frame
    content_text_frame.paragraphs[0].font.size = Pt(32)  # Set content font size

    if is_answer and answer:
        # For answers, add the answer text at the end and adjust the font size
        text_placeholder.text += f"\n\nAnswer: {answer}"
        content_text_frame.paragraphs[0].font.size = Pt(32)  # Adjust answer font size


# Loop through each question and add to the presentation
for i, row in df.iterrows():
    # Add a slide for the question with custom background color
    question_text = row['question']
    choices = f"A) {row['A)']}\nB) {row['B)']}\nC) {row['C)']}\nD) {row['D)']}"
    add_slide(prs, f"Question {i + 1}", question_text + "\n\n" + choices, bg_color=(255, 255, 255))  # Light blue background

    # Add a slide for the answer with a different background color
    correct_answer = row['correct_answer']
    if correct_answer in ['A', 'B', 'C', 'D']:
        answer_text = row[f'{correct_answer})']
        add_slide(prs, f"Answer to Question {i + 1}", question_text, answer=f"{correct_answer}: {answer_text}",
                  is_answer=True, bg_color=(255, 255, 255))  # Light orange background
    else:
        add_slide(prs, f"Answer to Question {i + 1}", question_text, answer=f"{correct_answer}", is_answer=True, bg_color=(255, 255, 255))


ppt_path = inputDir / 'quiz_presentation.pptx'
# Save the presentation to a file
prs.save(ppt_path)
